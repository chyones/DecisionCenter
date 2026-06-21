"""Exporter parity tests for ReportPolicy-driven section rendering."""

from __future__ import annotations

import io

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from reportlab.platypus import Paragraph, Table

from apps.edr.exporters.excel import to_excel
from apps.edr.exporters.pdf import _build_story, _make_doc
from apps.edr.exporters.powerpoint import to_powerpoint
from apps.edr.exporters.word import to_word


def _report(report_type: str) -> dict:
    return {
        "request_id": "r-export",
        "project_code": "PRJ-001",
        "project_identity": {"project_name": "Exporter Test Project", "project_code": "PRJ-001"},
        "query": "Show salary availability.",
        "report_type": report_type,
        "language": "en",
        "quality_gate_status": "passed",
        "executive_summary": [
            {"claim": "Salary records were checked.", "confidence": "medium", "evidence_ids": ["E1"]}
        ],
        "financial_snapshot": {
            "budget": {
                "value": 100,
                "currency": "AED",
                "evidence_id": "FIN-BUDGET",
                "status": "available",
            },
            "actual_cost": {
                "value": 80,
                "currency": "AED",
                "evidence_id": "FIN-ACTUAL",
                "status": "available",
            },
            "variance": {"value": 20, "currency": "AED", "formula": "budget - actual"},
        },
        "key_findings": [
            {"text": "Payroll extract exists.", "confidence": "high", "evidence_ids": ["E1"]}
        ],
        "root_causes": [
            {"text": "Root cause content must be suppressed.", "confidence": "low", "evidence_ids": ["E2"]}
        ],
        "delay_analysis": [
            {"text": "Delay content must be suppressed.", "confidence": "low", "evidence_ids": ["E3"]}
        ],
        "contractual_implications": [
            {
                "text": "Contract content must be suppressed.",
                "confidence": "low",
                "evidence_ids": ["E4"],
            }
        ],
        "recommended_actions": [
            {"text": "Share salary file.", "confidence": "medium", "evidence_ids": ["E1"]}
        ],
        "conflicts": [],
        "missing_data": [],
        "sources": [
            {
                "source_id": "E1",
                "source_type": "odoo",
                "title": "Payroll",
                "reference": "odoo://payroll",
                "date": "2026-06-21",
                "confidence": "high",
                "used_in": ["Key Findings"],
            }
        ],
        "connector_coverage": [],
    }


def _docx_headings(report: dict) -> list[str]:
    doc = Document(io.BytesIO(to_word(report)))
    return [
        paragraph.text
        for paragraph in doc.paragraphs
        if paragraph.style.name.startswith("Heading") and paragraph.text
    ]


def _story_text(report: dict) -> str:
    doc = _make_doc(io.BytesIO())
    story, _has_arabic = _build_story(report, doc)
    parts: list[str] = []

    def collect(value) -> None:
        if isinstance(value, Paragraph):
            parts.append(value.getPlainText())
        elif isinstance(value, Table):
            collect(value._cellvalues)
        elif isinstance(value, list | tuple):
            for item in value:
                collect(item)
        elif value is not None:
            parts.append(str(value))

    collect(story)
    return "\n".join(parts)


def _pptx_text(report: dict) -> str:
    prs = Presentation(io.BytesIO(to_powerpoint(report)))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    return "\n".join(parts)


def test_word_export_uses_policy_sections_and_contiguous_numbering() -> None:
    salary_headings = _docx_headings(_report("salary_payroll"))

    assert "2. Financial Snapshot — Odoo" not in salary_headings
    assert all("Root Causes" not in heading for heading in salary_headings)
    assert all("Delay Analysis" not in heading for heading in salary_headings)
    assert all("Contractual" not in heading for heading in salary_headings)
    assert "3. Recommended Actions — Proposal Only" in salary_headings
    assert salary_headings[-1] == "7. Quality Gate Status"

    full_headings = _docx_headings(_report("general_project_status"))
    assert "2. Financial Snapshot — Odoo" in full_headings
    assert full_headings[-1] == "11. Quality Gate Status"


def test_pdf_export_uses_policy_sections_and_contiguous_numbering() -> None:
    salary_text = _story_text(_report("salary_payroll"))

    assert "Financial Snapshot — Odoo" not in salary_text
    assert "Root Causes" not in salary_text
    assert "Delay Analysis" not in salary_text
    assert "Contractual / Commercial Implications" not in salary_text
    assert "3. Recommended Actions — Proposal Only" in salary_text
    assert "7. Quality Gate Status" in salary_text

    full_text = _story_text(_report("general_project_status"))
    assert "2. Financial Snapshot — Odoo" in full_text
    assert "11. Quality Gate Status" in full_text


def test_powerpoint_export_uses_policy_sections() -> None:
    salary_text = _pptx_text(_report("salary_payroll"))

    assert "Financial Snapshot — Odoo" not in salary_text
    assert "Root Causes & Delay Analysis" not in salary_text
    assert "All financial values sourced from Odoo ERP only." not in salary_text
    assert "Recommended Actions — Proposals Only" in salary_text

    full_text = _pptx_text(_report("general_project_status"))
    assert "Financial Snapshot — Odoo" in full_text
    assert "Root Causes & Delay Analysis" in full_text
    assert "All financial values sourced from Odoo ERP only." in full_text


def test_excel_export_uses_policy_sections() -> None:
    salary_wb = load_workbook(io.BytesIO(to_excel(_report("salary_payroll"))))
    assert "Financial Snapshot" not in salary_wb.sheetnames

    audit_sections = [
        cell.value
        for cell in salary_wb["Audit Data"]["A"]
        if cell.row > 1 and cell.value is not None
    ]
    assert "Key Findings" in audit_sections
    assert "Recommended Actions" in audit_sections
    assert "Root Causes" not in audit_sections
    assert "Delay Analysis" not in audit_sections
    assert "Contractual Implications" not in audit_sections

    full_wb = load_workbook(io.BytesIO(to_excel(_report("general_project_status"))))
    assert "Financial Snapshot" in full_wb.sheetnames
    full_audit_sections = [
        cell.value
        for cell in full_wb["Audit Data"]["A"]
        if cell.row > 1 and cell.value is not None
    ]
    assert "Root Causes" in full_audit_sections
    assert "Delay Analysis" in full_audit_sections
    assert "Contractual Implications" in full_audit_sections
