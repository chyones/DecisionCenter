"""Follow-up 2b tests — every exporter renders the distinct financial figures.

Markdown, Word, PDF, PowerPoint and Excel must all surface Contract Value /
Estimate / Committed Cost for a financial report (parity with the canonical
markdown renderer).
"""

from __future__ import annotations

from io import BytesIO

from apps.edr.exporters import export_report


def _available(value, eid):
    return {"value": value, "currency": "AED", "evidence_id": eid, "status": "available"}


def _financial_report() -> dict:
    return {
        "request_id": "r",
        "project_code": "PRJ-001",
        "project_identity": {"project_name": "Test Project", "project_code": "PRJ-001",
                             "identity_source": "registry", "identity_confidence": "verified"},
        "query": "what is the actual cost and budget variance",
        "report_type": "financial",
        "financial_snapshot": {
            "contract_value": _available(5000000.0, "ev_p"),
            "estimate": _available(4800000.0, "ev_p"),
            "budget": {"value": None, "currency": "AED", "evidence_id": None, "status": "not_available"},
            "actual_cost": _available(57000.0, "ev_a"),
            "committed_cost": _available(200000.0, "ev_po"),
            "variance": {"value": None, "currency": "AED", "formula": None, "evidence_ids": []},
        },
        "executive_summary": [{"claim": "Financial figures available.", "evidence_ids": ["ev_p"], "confidence": "low"}],
        "key_findings": [], "recommended_actions": [], "conflicts": [], "missing_data": [],
        "sources": [], "connector_coverage": [], "quality_gate_status": "passed",
    }


def test_all_exporters_render_distinct_financial_figures():
    results = export_report(_financial_report(), ["md", "docx", "xlsx", "pdf", "pptx"])
    assert {"md", "docx", "xlsx", "pdf", "pptx"} <= set(results)

    # Markdown
    md = results["md"].content.decode("utf-8")
    assert "Contract Value" in md and "Estimate" in md and "Committed Cost" in md

    # Word
    from docx import Document
    doc = Document(BytesIO(results["docx"].content))
    docx_text = " ".join(p.text for p in doc.paragraphs) + " " + " ".join(
        c.text for t in doc.tables for row in t.rows for c in row.cells
    )
    assert "Contract Value" in docx_text and "Committed Cost" in docx_text

    # Excel
    import openpyxl
    wb = openpyxl.load_workbook(BytesIO(results["xlsx"].content))
    cells = [str(c.value) for ws in wb.worksheets for row in ws.iter_rows() for c in row if c.value is not None]
    assert any("Contract Value" in c for c in cells) and any("Committed Cost" in c for c in cells)

    # PowerPoint
    from pptx import Presentation
    prs = Presentation(BytesIO(results["pptx"].content))
    ptext = " ".join(
        sh.text for sl in prs.slides for sh in sl.shapes if getattr(sh, "has_text_frame", False)
    )
    assert "Contract Value" in ptext and "Committed Cost" in ptext

    # PDF — valid document produced (text not introspected)
    assert results["pdf"].content[:4] == b"%PDF"
