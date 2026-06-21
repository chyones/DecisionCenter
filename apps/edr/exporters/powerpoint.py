"""PowerPoint (.pptx) exporter — for executive presentation summaries."""

from __future__ import annotations

import io
from datetime import datetime, timezone

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from apps.edr.graph.report_policy import (
    SEC_DELAY_ANALYSIS,
    SEC_FINANCIAL_SNAPSHOT,
    SEC_ROOT_CAUSES,
    policy_for,
)

_BRAND_BLUE = RGBColor(0x1F, 0x38, 0x64)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_ACCENT = RGBColor(0xCC, 0xCC, 0xFF)
_DARK_GREY = RGBColor(0x80, 0x80, 0x80)


def to_powerpoint(report: dict) -> bytes:
    """Generate an executive summary PowerPoint from the canonical JSON report dict."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.63)  # 16:9

    request_id = report.get("request_id", "N/A")
    project_code = report.get("project_code") or "N/A"
    query = report.get("query") or report.get("question", "N/A")
    qg_status = report.get("quality_gate_status", "not_run")
    generated_at = datetime.now(timezone.utc).isoformat(timespec="seconds")
    policy = policy_for(report.get("report_type", "executive_decision"))

    _title_slide(prs, query, request_id, project_code)
    _summary_slide(prs, report)
    if policy.renders(SEC_FINANCIAL_SNAPSHOT):
        _financial_slide(prs, report)
    _findings_slide(prs, report)
    if policy.renders(SEC_ROOT_CAUSES) or policy.renders(SEC_DELAY_ANALYSIS):
        _causes_delay_slide(prs, report, include_root=policy.renders(SEC_ROOT_CAUSES), include_delay=policy.renders(SEC_DELAY_ANALYSIS))
    _actions_slide(prs, report)
    _conflicts_missing_slide(prs, report)
    _status_slide(
        prs,
        request_id,
        project_code,
        qg_status,
        generated_at,
        include_financial=policy.renders(SEC_FINANCIAL_SNAPSHOT),
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Slide builders ────────────────────────────────────────────────────────────


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def _content_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[1])  # title + content


def _title_slide(prs: Presentation, query: str, request_id: str, project_code: str) -> None:
    slide = _blank_slide(prs)

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _BRAND_BLUE

    def _tb(left, top, w, h) -> object:
        return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))

    title_tf = _tb(0.8, 1.4, 8.4, 1.4).text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.text = "Executive Decision Report"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = _WHITE
    p.alignment = PP_ALIGN.CENTER

    sub_tf = _tb(0.8, 3.0, 8.4, 1.2).text_frame
    sub_tf.word_wrap = True
    p2 = sub_tf.paragraphs[0]
    p2.text = query[:120] + ("…" if len(query) > 120 else "")
    p2.font.size = Pt(13)
    p2.font.color.rgb = _WHITE
    p2.alignment = PP_ALIGN.CENTER

    meta_tf = _tb(0.8, 4.3, 8.4, 0.7).text_frame
    p3 = meta_tf.paragraphs[0]
    p3.text = f"Project: {project_code}  |  ID: {request_id}"
    p3.font.size = Pt(10)
    p3.font.color.rgb = _ACCENT
    p3.alignment = PP_ALIGN.CENTER


def _summary_slide(prs: Presentation, report: dict) -> None:
    slide = _content_slide(prs)
    _set_slide_title(slide, "Executive Summary")
    summary = report.get("executive_summary", [])
    bullets: list[str] = []
    if isinstance(summary, str):
        bullets = [summary]
    else:
        for item in summary:
            if isinstance(item, dict):
                bullets.append(f"{item.get('claim', '')} [{item.get('confidence', 'medium')}]")
    _fill_content(slide, bullets or ["No summary available."])


def _financial_slide(prs: Presentation, report: dict) -> None:
    slide = _content_slide(prs)
    _set_slide_title(slide, "Financial Snapshot — Odoo")
    fs = report.get("financial_snapshot") or {}
    bullets: list[str] = []
    if isinstance(fs, dict):
        budget = fs.get("budget") or {}
        actual = fs.get("actual_cost") or {}
        variance = fs.get("variance") or {}
        currency = (budget.get("currency") if isinstance(budget, dict) else None) or "AED"
        bullets.append(_fin_bullet("Budget", budget, currency))
        bullets.append(_fin_bullet("Actual Cost", actual, currency))
        if isinstance(variance, dict):
            v = variance.get("value")
            c = variance.get("currency", currency)
            formula = variance.get("formula", "")
            val_str = f"{v:,.2f} {c}" if v is not None else "Not available"
            bullets.append(f"Variance: {val_str}{' (' + formula + ')' if formula else ''}")
        bullets.append("Source: Odoo ERP (authoritative)")
    _fill_content(slide, bullets or ["Financial data not available."])


def _findings_slide(prs: Presentation, report: dict) -> None:
    slide = _content_slide(prs)
    _set_slide_title(slide, "Key Findings")
    items = report.get("key_findings", [])
    bullets = [
        f"{i.get('text', '')} [{i.get('confidence', 'medium')}]"
        if isinstance(i, dict) else str(i)
        for i in items
    ]
    _fill_content(slide, bullets or ["No findings available."])


def _causes_delay_slide(
    prs: Presentation, report: dict, *, include_root: bool, include_delay: bool
) -> None:
    slide = _content_slide(prs)
    _set_slide_title(slide, "Root Causes & Delay Analysis")
    bullets: list[str] = []
    if include_root:
        for item in report.get("root_causes", []):
            text = item.get("text", str(item)) if isinstance(item, dict) else str(item)
            bullets.append(f"Root Cause: {text}")
    if include_delay:
        for item in report.get("delay_analysis", []):
            text = item.get("text", str(item)) if isinstance(item, dict) else str(item)
            bullets.append(f"Delay: {text}")
    _fill_content(slide, bullets or ["No data available."])


def _actions_slide(prs: Presentation, report: dict) -> None:
    slide = _content_slide(prs)
    _set_slide_title(slide, "Recommended Actions — Proposals Only")
    actions = report.get("recommended_actions", [])
    bullets = [
        f"→ {i.get('text', '')} [{i.get('confidence', 'medium')}]"
        if isinstance(i, dict) else f"→ {i}"
        for i in actions
    ]
    _fill_content(slide, bullets or ["No recommended actions."])


def _conflicts_missing_slide(prs: Presentation, report: dict) -> None:
    slide = _content_slide(prs)
    _set_slide_title(slide, "Conflicts & Missing Data")
    bullets: list[str] = []
    for c in report.get("conflicts", []):
        if isinstance(c, dict):
            desc = c.get("description", "")[:80]
            bullets.append(f"⚠ Conflict ({c.get('conflict_type', '?')}): {desc}")
    for m in report.get("missing_data", []):
        bullets.append(f"⚠ Missing: {m}")
    _fill_content(slide, bullets or ["No conflicts or missing data."])


def _status_slide(
    prs: Presentation,
    request_id: str,
    project_code: str,
    qg_status: str,
    generated_at: str,
    *,
    include_financial: bool,
) -> None:
    slide = _content_slide(prs)
    _set_slide_title(slide, "Report Status")
    bullets = [
        f"Quality Gate: {qg_status}",
        f"Request ID: {request_id}",
        f"Project: {project_code}",
        f"Generated: {generated_at}",
        "This presentation is for executive review only.",
    ]
    if include_financial:
        bullets.insert(4, "All financial values sourced from Odoo ERP only.")
    _fill_content(slide, bullets)


# ── Helpers ───────────────────────────────────────────────────────────────────


def _set_slide_title(slide, text: str) -> None:
    title_shape = slide.shapes.title
    title_shape.text = text
    p = title_shape.text_frame.paragraphs[0]
    p.font.color.rgb = _BRAND_BLUE
    p.font.size = Pt(22)
    p.font.bold = True


def _fill_content(slide, bullets: list[str]) -> None:
    content_ph = slide.placeholders[1]
    tf = content_ph.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, bullet in enumerate(bullets[:8]):  # cap at 8 per slide
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(13)
        p.level = 0


def _fin_bullet(label: str, node: dict, currency: str) -> str:
    if not isinstance(node, dict):
        return f"{label}: Not available"
    v = node.get("value")
    c = node.get("currency", currency)
    return f"{label}: {v:,.2f} {c}" if v is not None else f"{label}: Not available"
